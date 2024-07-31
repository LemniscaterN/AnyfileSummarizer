import argparse
import os
import subprocess

def extract_text_from_pdf(pdf_path):
    import fitz
    document = fitz.open(pdf_path)
    text = []
    for page_num in range(len(document)):
        page = document.load_page(page_num)
        text.append(page.get_text())
    return "\n".join(text)

parser = argparse.ArgumentParser()
parser.add_argument('-pdf', help='.pdf file path is required.')
parser.add_argument('-doc', help='.doc file or .docx path is required.')
parser.add_argument('-ptx', help='.ptx file or .pptx path is required.')
args = parser.parse_args()

target_text = ""

if args.pdf and os.path.exists(args.pdf):
    target_text = extract_text_from_pdf(args.pdf)
elif args.doc and os.path.exists(args.doc):
    from docx import Document
    doc = Document(args.doc)
    for paragraph in doc.paragraphs:
        if(len(paragraph.text)!=0):
            target_text+=f"# {paragraph.text}"
elif args.ptx and os.path.exists(args.ptx):
    from pptx import Presentation
    import re
    prs = Presentation(args.ptx)
    filename = os.path.splitext(os.path.basename(args.ptx))[0]
    for i,slide in  enumerate(prs.slides):
        target_text += "# "
        for shape in slide.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame is not None:
                target_text += shape.text_frame.text
        if(len(target_text)>1000):
            target_text = re.sub(r'\s+', ' ',target_text).replace("\t", " ").replace("\n", " ")
            target_text = subprocess.run(['ollama', 'run', 'elyza:jp8b',f'Summarize please. # {filename}{target_text}'],capture_output=True, text=True).stdout
    target_text = re.sub(r'\s+', ' ',target_text).replace("\t", " ").replace("\n", " ")

if(target_text != ""):
    print(f"要約結果：\n{subprocess.run(['ollama', 'run', 'elyza:jp8b',f'Summarize please. # {target_text}'],capture_output=True, text=True).stdout}")